"""PPTX parser with slide visuals and multi-modal support"""

import logging
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from core.parsers.interfaces import (
    BaseParser,
    Document,
    DocumentMetadata,
    DocumentType,
    ParseError,
    Segment,
    SegmentType,
    TextSubtype,
    VisualSubtype,
    MetadataSubtype,
    VisualElement,
    VisualElementType,
)

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """
    PPTX parser with slide visuals and multi-modal support
    
    Features:
    - Slide-based text extraction
    - Image and diagram extraction
    - Chart analysis from slides
    - Multi-modal content with VLM descriptions
    - Slide-level segmentation
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None, enable_vlm: bool = True):
        """Initialize PPTX parser"""
        super().__init__(config, enable_vlm)
        self.supported_extensions = {".pptx", ".ppt"}
        
        # PPTX-specific configuration
        self.config.setdefault("extract_images", True)
        self.config.setdefault("extract_charts", True)
        self.config.setdefault("extract_shapes", True)
        self.config.setdefault("include_slide_notes", True)
        self.config.setdefault("include_slide_titles", True)
        self.config.setdefault("image_min_size", 50)
        
        logger.info(f"Initialized PPTX parser with VLM: {enable_vlm}")
    
    def can_parse(self, file_path: Path) -> bool:
        """Check if file is a PPTX presentation"""
        return file_path.suffix.lower() in self.supported_extensions
    
    async def parse(self, file_path: Path) -> Document:
        """
        Parse PPTX presentation with slide visuals
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Document with text segments and visual elements
            
        Raises:
            ParseError: If parsing fails
        """
        self.validate_file(file_path)
        
        try:
            logger.info(f"Starting PPTX parsing: {file_path.name}")
            
            # Load presentation
            presentation = Presentation(str(file_path))
            
            # Extract metadata
            metadata = self._extract_pptx_metadata(file_path, presentation)
            
            # Extract visual elements
            visual_elements = []
            if self.config.get("extract_images", True) or self.config.get("extract_charts", True):
                visual_elements = self._extract_visual_elements(file_path, presentation)
            
            # Analyze visual elements with VLM
            if visual_elements and self.enable_vlm:
                document_context = {
                    "document_title": metadata.title,
                    "document_type": "PPTX",
                    "slide_count": len(presentation.slides)
                }
                visual_elements = await self.analyze_visual_elements(visual_elements, document_context)
            
            # Create segments from slides
            segments = self._create_segments(presentation, visual_elements)
            
            # Create visual segments
            visual_segments = self._create_visual_segments(visual_elements)
            
            # Merge all segments maintaining slide order
            all_segments = self._merge_segments_by_slide(segments, visual_segments)
            
            # Build full document content
            full_content = self._build_document_content(all_segments, visual_elements)
            
            # Create document
            document = self.create_document(
                file_path,
                segments=all_segments,
                metadata=metadata,
                visual_elements=visual_elements,
                content=full_content,
                raw_data=presentation
            )
            
            logger.info(f"PPTX parsing completed: {document.total_segments} segments, "
                       f"{document.total_visual_elements} visual elements")
            
            return document
            
        except Exception as e:
            logger.error(f"PPTX parsing failed: {e}", exc_info=True)
            raise ParseError(f"Failed to parse PPTX {file_path.name}: {str(e)}")
    
    def _extract_pptx_metadata(self, file_path: Path, presentation: Presentation) -> DocumentMetadata:
        """Extract metadata from PPTX presentation"""
        try:
            # Get base metadata
            metadata = self.extract_metadata(file_path)
            metadata.document_type = DocumentType.PPTX
            
            # Extract presentation properties
            core_props = presentation.core_properties
            if core_props:
                metadata.title = core_props.title or metadata.title
                metadata.author = core_props.author
                metadata.created = core_props.created
                metadata.modified = core_props.modified
            
            # Count slides as pages
            metadata.page_count = len(presentation.slides)
            
            # Add custom metadata
            metadata.custom_metadata.update({
                "slides": len(presentation.slides),
                "slide_layouts": len(presentation.slide_layouts),
                "slide_masters": len(presentation.slide_masters),
                "has_images": self._has_images(presentation),
                "has_charts": self._has_charts(presentation),
                "format": "PPTX"
            })
            
            return metadata
            
        except Exception as e:
            logger.warning(f"Failed to extract PPTX metadata: {e}")
            metadata = self.extract_metadata(file_path)
            metadata.document_type = DocumentType.PPTX
            return metadata
    
    def _has_images(self, presentation: Presentation) -> bool:
        """Check if presentation contains images"""
        try:
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        return True
            return False
        except Exception:
            return False
    
    def _has_charts(self, presentation: Presentation) -> bool:
        """Check if presentation contains charts"""
        try:
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        return True
            return False
        except Exception:
            return False
    
    def _extract_visual_elements(self, file_path: Path, presentation: Presentation) -> List[VisualElement]:
        """Extract visual elements from PPTX presentation"""
        visual_elements = []
        
        try:
            # Extract from each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                
                # Extract visual elements from slide shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Handle different shape types
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.config.get("extract_images", True):
                            visual_element = self._extract_image_shape(shape, slide_number, shape_idx, file_path)
                            if visual_element:
                                visual_elements.append(visual_element)
                        
                        elif shape.shape_type == MSO_SHAPE_TYPE.CHART and self.config.get("extract_charts", True):
                            visual_element = self._extract_chart_shape(shape, slide_number, shape_idx, file_path)
                            if visual_element:
                                visual_elements.append(visual_element)
                        
                        elif self.config.get("extract_shapes", True) and shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP]:
                            # Extract other shapes that might contain visual information
                            visual_element = self._extract_shape_element(shape, slide_number, shape_idx, file_path)
                            if visual_element:
                                visual_elements.append(visual_element)
                                
                    except Exception as e:
                        logger.warning(f"Failed to extract shape {shape_idx} from slide {slide_number}: {e}")
                        continue
                
                # Extract images from slide media relations
                slide_visual_elements = self._extract_slide_media(file_path, slide_number)
                visual_elements.extend(slide_visual_elements)
            
            # Remove duplicates
            unique_elements = {}
            for element in visual_elements:
                if element.content_hash not in unique_elements:
                    unique_elements[element.content_hash] = element
            
            visual_elements = list(unique_elements.values())
            
            logger.info(f"Extracted {len(visual_elements)} visual elements from PPTX")
            return visual_elements
            
        except Exception as e:
            logger.error(f"Visual element extraction failed: {e}")
            return []
    
    def _extract_image_shape(self, shape, slide_number: int, shape_idx: int, file_path: Path) -> Optional[VisualElement]:
        """Extract image from a picture shape"""
        try:
            # Get image data
            image_data = self._get_shape_image_data(shape)
            if not image_data or len(image_data) < 1000:
                return None
            
            # Determine image type
            img_type = self._determine_image_type(image_data)
            
            visual_element = VisualElement(
                element_type=img_type,
                source_format=DocumentType.PPTX,
                content_hash=VisualElement.create_hash(image_data),
                raw_data=image_data,
                page_or_slide=slide_number,
                segment_reference=f"slide_{slide_number}",
                file_extension=self._get_image_extension(image_data),
                bounding_box=self._get_shape_bounds(shape),
                analysis_metadata={
                    "slide_number": slide_number,
                    "shape_index": shape_idx,
                    "shape_type": "picture",
                    "extraction_method": "shape_picture"
                }
            )
            
            return visual_element
            
        except Exception as e:
            logger.debug(f"Failed to extract image from shape: {e}")
            return None
    
    def _extract_chart_shape(self, shape, slide_number: int, shape_idx: int, file_path: Path) -> Optional[VisualElement]:
        """Extract chart from a chart shape"""
        try:
            # Analyze chart data
            chart_data = self._analyze_chart_shape(shape)
            
            # Create text representation of chart
            chart_text = self._chart_data_to_text(chart_data)
            chart_bytes = chart_text.encode('utf-8')
            
            visual_element = VisualElement(
                element_type=VisualElementType.CHART,
                source_format=DocumentType.PPTX,
                content_hash=VisualElement.create_hash(chart_bytes),
                raw_data=chart_bytes,
                page_or_slide=slide_number,
                segment_reference=f"slide_{slide_number}",
                file_extension="txt",
                bounding_box=self._get_shape_bounds(shape),
                extracted_data=chart_data,
                analysis_metadata={
                    "slide_number": slide_number,
                    "shape_index": shape_idx,
                    "shape_type": "chart",
                    "extraction_method": "shape_chart"
                }
            )
            
            return visual_element
            
        except Exception as e:
            logger.debug(f"Failed to extract chart from shape: {e}")
            return None
    
    def _extract_shape_element(self, shape, slide_number: int, shape_idx: int, file_path: Path) -> Optional[VisualElement]:
        """Extract visual information from other shapes"""
        try:
            # Get shape text and properties
            shape_info = self._analyze_shape(shape)
            
            if not shape_info.get("has_visual_content", False):
                return None
            
            # Create text representation
            shape_text = self._shape_info_to_text(shape_info)
            shape_bytes = shape_text.encode('utf-8')
            
            visual_element = VisualElement(
                element_type=VisualElementType.DIAGRAM,
                source_format=DocumentType.PPTX,
                content_hash=VisualElement.create_hash(shape_bytes),
                raw_data=shape_bytes,
                page_or_slide=slide_number,
                segment_reference=f"slide_{slide_number}",
                file_extension="txt",
                bounding_box=self._get_shape_bounds(shape),
                extracted_data=shape_info,
                analysis_metadata={
                    "slide_number": slide_number,
                    "shape_index": shape_idx,
                    "shape_type": str(shape.shape_type),
                    "extraction_method": "shape_analysis"
                }
            )
            
            return visual_element
            
        except Exception as e:
            logger.debug(f"Failed to extract shape element: {e}")
            return None
    
    def _extract_slide_media(self, file_path: Path, slide_number: int) -> List[VisualElement]:
        """Extract media files from slide relations"""
        visual_elements = []
        
        try:
            # Open PPTX as zip to access media files
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Look for media files in the presentation
                media_files = [f for f in zip_file.namelist() if f.startswith('ppt/media/')]
                
                for media_file in media_files:
                    try:
                        # Read media file
                        media_data = zip_file.read(media_file)
                        
                        # Skip small files
                        if len(media_data) < 1000:
                            continue
                        
                        # Determine if it's an image
                        if self._is_image_file(media_file):
                            img_type = self._determine_image_type(media_data)
                            
                            visual_element = VisualElement(
                                element_type=img_type,
                                source_format=DocumentType.PPTX,
                                content_hash=VisualElement.create_hash(media_data),
                                raw_data=media_data,
                                page_or_slide=slide_number,
                                segment_reference=f"slide_{slide_number}",
                                file_extension=Path(media_file).suffix.lower().lstrip('.'),
                                analysis_metadata={
                                    "media_file": media_file,
                                    "slide_number": slide_number,
                                    "extraction_method": "media_file"
                                }
                            )
                            
                            visual_elements.append(visual_element)
                            
                    except Exception as e:
                        logger.warning(f"Failed to extract media file {media_file}: {e}")
                        continue
        
        except Exception as e:
            logger.error(f"Failed to extract slide media: {e}")
        
        return visual_elements
    
    def _get_shape_image_data(self, shape) -> Optional[bytes]:
        """Get image data from a picture shape"""
        try:
            if hasattr(shape, 'image') and shape.image:
                return shape.image.blob
            return None
        except Exception as e:
            logger.debug(f"Failed to get shape image data: {e}")
            return None
    
    def _get_shape_bounds(self, shape) -> Optional[Dict[str, float]]:
        """Get bounding box for shape"""
        try:
            return {
                "left": float(shape.left),
                "top": float(shape.top),
                "width": float(shape.width),
                "height": float(shape.height)
            }
        except Exception:
            return None
    
    def _analyze_chart_shape(self, shape) -> Dict[str, Any]:
        """Analyze chart shape and extract data"""
        try:
            chart_data = {
                "chart_type": "unknown",
                "title": None,
                "has_data": False,
                "series_count": 0
            }
            
            if hasattr(shape, 'chart') and shape.chart:
                chart = shape.chart
                
                # Get chart title
                if hasattr(chart, 'chart_title') and chart.chart_title:
                    chart_data["title"] = chart.chart_title.text_frame.text
                
                # Get chart type
                if hasattr(chart, 'chart_type'):
                    chart_data["chart_type"] = str(chart.chart_type)
                
                # Get series information
                if hasattr(chart, 'series'):
                    chart_data["series_count"] = len(chart.series)
                    chart_data["has_data"] = chart_data["series_count"] > 0
            
            return chart_data
            
        except Exception as e:
            logger.debug(f"Chart analysis failed: {e}")
            return {"chart_type": "unknown", "error": str(e)}
    
    def _analyze_shape(self, shape) -> Dict[str, Any]:
        """Analyze shape for visual content"""
        try:
            shape_info = {
                "shape_type": str(shape.shape_type),
                "has_text": False,
                "has_visual_content": False,
                "text_content": None,
                "shape_name": None
            }
            
            # Check for text content
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shape_info["has_text"] = True
                    shape_info["text_content"] = text
            
            # Check for shape name
            if hasattr(shape, 'name'):
                shape_info["shape_name"] = shape.name
            
            # Determine if shape has visual content
            shape_info["has_visual_content"] = (
                shape_info["has_text"] or 
                shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP]
            )
            
            return shape_info
            
        except Exception as e:
            logger.debug(f"Shape analysis failed: {e}")
            return {"shape_type": "unknown", "has_visual_content": False}
    
    def _chart_data_to_text(self, chart_data: Dict[str, Any]) -> str:
        """Convert chart data to text representation"""
        try:
            text_parts = []
            
            if chart_data.get("title"):
                text_parts.append(f"Chart Title: {chart_data['title']}")
            
            text_parts.append(f"Chart Type: {chart_data.get('chart_type', 'Unknown')}")
            
            if chart_data.get("has_data"):
                text_parts.append(f"Series Count: {chart_data.get('series_count', 0)}")
            else:
                text_parts.append("No data series found")
            
            return "\n".join(text_parts)
            
        except Exception:
            return f"Chart: {chart_data.get('chart_type', 'Unknown')}"
    
    def _shape_info_to_text(self, shape_info: Dict[str, Any]) -> str:
        """Convert shape info to text representation"""
        try:
            text_parts = []
            
            text_parts.append(f"Shape Type: {shape_info.get('shape_type', 'Unknown')}")
            
            if shape_info.get("shape_name"):
                text_parts.append(f"Shape Name: {shape_info['shape_name']}")
            
            if shape_info.get("has_text") and shape_info.get("text_content"):
                text_parts.append(f"Text Content: {shape_info['text_content']}")
            
            return "\n".join(text_parts)
            
        except Exception:
            return f"Shape: {shape_info.get('shape_type', 'Unknown')}"
    
    def _is_image_file(self, filename: str) -> bool:
        """Check if file is an image"""
        image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
        return Path(filename).suffix.lower() in image_extensions
    
    def _determine_image_type(self, image_data: bytes) -> VisualElementType:
        """Determine visual element type for image"""
        try:
            from PIL import Image
            import io
            
            img = Image.open(io.BytesIO(image_data))
            width, height = img.size
            
            # Heuristics for image classification
            aspect_ratio = width / height
            
            if aspect_ratio > 2.0:
                return VisualElementType.CHART
            elif 0.8 <= aspect_ratio <= 1.2:
                return VisualElementType.DIAGRAM
            else:
                return VisualElementType.IMAGE
                
        except Exception:
            return VisualElementType.UNKNOWN_VISUAL
    
    def _get_image_extension(self, image_data: bytes) -> str:
        """Get file extension for image"""
        try:
            from PIL import Image
            import io
            
            img = Image.open(io.BytesIO(image_data))
            format_map = {
                'JPEG': 'jpg',
                'PNG': 'png',
                'BMP': 'bmp',
                'GIF': 'gif',
                'TIFF': 'tiff'
            }
            return format_map.get(img.format, 'png')
        except Exception:
            return 'png'
    
    def _create_segments(self, presentation: Presentation, visual_elements: List[VisualElement]) -> List[Segment]:
        """Create segments from presentation slides"""
        segments = []
        
        try:
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                
                # Extract slide content
                slide_content = self._extract_slide_content(slide, slide_number)
                
                if slide_content:
                    # Find visual references for this slide
                    visual_refs = []
                    for ve in visual_elements:
                        if ve.page_or_slide == slide_number:
                            visual_refs.append(ve.content_hash)
                    
                    # Get slide title
                    slide_title = self._get_slide_title(slide)
                    
                    segment = Segment(
                        content=slide_content,
                        page_number=slide_number,
                        segment_index=slide_idx,
                        segment_type=SegmentType.METADATA,
                        segment_subtype=MetadataSubtype.SLIDE.value,
                        visual_references=visual_refs,
                        metadata={
                            "slide_number": slide_number,
                            "slide_index": slide_idx,
                            "slide_title": slide_title if slide_title else "No title",
                            "shape_count": len(slide.shapes),
                            "has_title": self._slide_has_title(slide),
                            "has_images": any(ve.element_type == VisualElementType.IMAGE for ve in visual_elements if ve.page_or_slide == slide_number),
                            "has_charts": any(ve.element_type == VisualElementType.CHART for ve in visual_elements if ve.page_or_slide == slide_number)
                        }
                    )
                    segments.append(segment)
                
                # Extract slide notes if enabled
                if self.config.get("include_slide_notes", True):
                    notes_segment = self._extract_slide_notes(slide, slide_number, len(segments))
                    if notes_segment:
                        segments.append(notes_segment)
            
            logger.info(f"Created {len(segments)} segments from PPTX")
            return segments
            
        except Exception as e:
            logger.error(f"Segment creation failed: {e}")
            return []
    
    def _extract_slide_content(self, slide, slide_number: int) -> str:
        """Extract text content from slide"""
        try:
            content_parts = []
            title_text = None
            
            # First, try to find the title placeholder specifically
            if hasattr(slide, 'placeholders'):
                for placeholder in slide.placeholders:
                    try:
                        # Type 0 = Title, Type 1 = Center Title
                        if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type in [0, 1]:
                            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                                title_text = placeholder.text_frame.text.strip()
                                if title_text:
                                    content_parts.append(f"TITLE: {title_text}")
                                    break
                    except Exception:
                        pass
            
            # If no title found in placeholders, check shapes
            if not title_text:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            # Check if it's a title shape
                            if self._is_title_shape(shape) and not title_text:
                                content_parts.append(f"TITLE: {text}")
                                title_text = text
                            elif text != title_text:  # Avoid duplicating title text
                                content_parts.append(text)
                    elif hasattr(shape, 'table') and shape.table:
                        # Handle table content
                        table_text = self._extract_table_content(shape.table)
                        if table_text:
                            content_parts.append(f"TABLE:\n{table_text}")
            else:
                # Process remaining shapes (excluding the title we already found)
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text and text != title_text and not self._is_title_shape(shape):
                            content_parts.append(text)
                    elif hasattr(shape, 'table') and shape.table:
                        # Handle table content
                        table_text = self._extract_table_content(shape.table)
                        if table_text:
                            content_parts.append(f"TABLE:\n{table_text}")
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            logger.debug(f"Failed to extract slide content: {e}")
            return ""
    
    def _extract_slide_notes(self, slide, slide_number: int, segment_idx: int) -> Optional[Segment]:
        """Extract notes from slide"""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text = ""
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            notes_text += text + "\n"
                
                if notes_text.strip():
                    return Segment(
                        content=notes_text.strip(),
                        page_number=slide_number,
                        segment_index=segment_idx,
                        segment_type=SegmentType.TEXT,
                        segment_subtype=TextSubtype.FOOTNOTE.value,
                        metadata={
                            "slide_number": slide_number,
                            "content_type": "notes"
                        }
                    )
            
            return None
            
        except Exception as e:
            logger.debug(f"Failed to extract slide notes: {e}")
            return None
    
    def _extract_table_content(self, table) -> str:
        """Extract content from table"""
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    cells.append(cell_text)
                if any(cells):
                    rows.append(" | ".join(cells))
            
            return "\n".join(rows)
            
        except Exception as e:
            logger.debug(f"Failed to extract table content: {e}")
            return ""
    
    def _slide_has_title(self, slide) -> bool:
        """Check if slide has a title"""
        try:
            # First check placeholders
            if hasattr(slide, 'placeholders'):
                for placeholder in slide.placeholders:
                    try:
                        if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type in [0, 1]:
                            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                                if placeholder.text_frame.text.strip():
                                    return True
                    except Exception:
                        pass
            
            # Then check shapes
            for shape in slide.shapes:
                if self._is_title_shape(shape):
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        if shape.text_frame.text.strip():
                            return True
            return False
        except Exception:
            return False
    
    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract slide title if present"""
        try:
            # First check placeholders
            if hasattr(slide, 'placeholders'):
                for placeholder in slide.placeholders:
                    try:
                        if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type in [0, 1]:
                            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                                title = placeholder.text_frame.text.strip()
                                if title:
                                    return title
                    except Exception:
                        pass
            
            # Then check shapes
            for shape in slide.shapes:
                if self._is_title_shape(shape):
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        title = shape.text_frame.text.strip()
                        if title:
                            return title
            return None
        except Exception:
            return None
    
    def _is_title_shape(self, shape) -> bool:
        """Check if shape is a title"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                # Type 0 = Title, Type 1 = Center Title
                return shape.placeholder_format.type in [0, 1]
            
            # Also check shape name for common title patterns
            if hasattr(shape, 'name') and shape.name:
                name_lower = shape.name.lower()
                if any(title_word in name_lower for title_word in ['title', 'titel', 'heading']):
                    return True
            
            return False
        except Exception:
            return False
    
    def _build_document_content(self, segments: List[Segment], visual_elements: List[VisualElement]) -> str:
        """Build full document content"""
        content_parts = []
        
        # Add text content
        for segment in segments:
            if hasattr(segment.segment_type, 'value') and segment.segment_type == SegmentType.METADATA and segment.segment_subtype == MetadataSubtype.SLIDE.value:
                content_parts.append(f"## Slide {segment.page_number}\n{segment.content}")
            elif hasattr(segment.segment_type, 'value') and segment.segment_type == SegmentType.TEXT and segment.segment_subtype == TextSubtype.FOOTNOTE.value:
                content_parts.append(f"### Notes for Slide {segment.page_number}\n{segment.content}")
            else:
                content_parts.append(segment.content)
        
        # Add visual element descriptions
        for visual in visual_elements:
            if visual.vlm_description:
                visual_desc = f"\n[VISUAL ELEMENT - {visual.element_type.value.upper()}]"
                if visual.page_or_slide:
                    visual_desc += f" (Slide {visual.page_or_slide})"
                visual_desc += f": {visual.vlm_description}"
                content_parts.append(visual_desc)
        
        return "\n\n".join(content_parts)
    
    def _create_visual_segments(self, visual_elements: List[VisualElement]) -> List[Segment]:
        """Create segments from visual elements"""
        visual_segments = []
        
        for visual_elem in visual_elements:
            # Determine visual subtype
            if visual_elem.element_type == VisualElementType.IMAGE:
                subtype = VisualSubtype.IMAGE.value
            elif visual_elem.element_type == VisualElementType.CHART:
                subtype = VisualSubtype.CHART.value
            elif visual_elem.element_type == VisualElementType.DIAGRAM:
                subtype = VisualSubtype.DIAGRAM.value
            elif visual_elem.element_type == VisualElementType.FORMULA:
                subtype = VisualSubtype.FORMULA.value
            else:
                subtype = VisualSubtype.IMAGE.value  # Default
            
            # Create placeholder content
            content = f"[{visual_elem.element_type.value.upper()}: Placeholder]"
            
            # Create visual segment
            segment = Segment(
                content=content,
                page_number=visual_elem.page_or_slide,  # slide number
                segment_type=SegmentType.VISUAL,
                segment_subtype=subtype,
                visual_references=[visual_elem.content_hash],
                metadata={
                    "visual_type": visual_elem.element_type.value,
                    "slide_number": visual_elem.page_or_slide,
                    "extraction_method": visual_elem.analysis_metadata.get("extraction_method", ""),
                    "source": "pptx_visual_extraction"
                }
            )
            
            visual_segments.append(segment)
        
        return visual_segments
    
    def _merge_segments_by_slide(self, text_segments: List[Segment], visual_segments: List[Segment]) -> List[Segment]:
        """Merge text and visual segments, maintaining slide order"""
        # Group segments by slide
        segments_by_slide = {}
        
        # Add text segments
        for segment in text_segments:
            slide = segment.page_number or 1
            if slide not in segments_by_slide:
                segments_by_slide[slide] = []
            segments_by_slide[slide].append(segment)
        
        # Add visual segments
        for segment in visual_segments:
            slide = segment.page_number or 1
            if slide not in segments_by_slide:
                segments_by_slide[slide] = []
            segments_by_slide[slide].append(segment)
        
        # Merge and re-index
        all_segments = []
        for slide in sorted(segments_by_slide.keys()):
            slide_segments = segments_by_slide[slide]
            # Sort: slide metadata first, then visuals, then notes
            slide_segments.sort(key=lambda s: (
                0 if s.segment_type == SegmentType.METADATA and s.segment_subtype == MetadataSubtype.SLIDE.value else
                1 if s.segment_type == SegmentType.VISUAL else
                2
            ))
            
            for segment in slide_segments:
                segment.segment_index = len(all_segments)
                all_segments.append(segment)
        
        return all_segments