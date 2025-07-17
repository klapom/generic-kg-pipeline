#!/usr/bin/env python3
"""
Detailed test of Office parsers to verify content extraction accuracy
"""

import sys
import os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../../..')))

import asyncio
import logging
from pathlib import Path
from datetime import datetime
import json
from typing import Any, Dict
import subprocess

# Set up logging
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
log_filename = f"data/output/Office_Parsers_detailed_{timestamp}.log"

# Create formatter
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')

# Console handler
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(formatter)

# File handler
os.makedirs("data/output", exist_ok=True)
file_handler = logging.FileHandler(log_filename, mode='w', encoding='utf-8')
file_handler.setLevel(logging.DEBUG)
file_handler.setFormatter(formatter)

# Set up root logger
logging.basicConfig(
    level=logging.DEBUG,
    handlers=[console_handler, file_handler]
)
logger = logging.getLogger(__name__)


def save_detailed_output(filename: str, document, parser_type: str):
    """Save detailed document content to file for comparison"""
    output = {
        "parser": parser_type,
        "file": filename,
        "metadata": {
            "title": document.metadata.title,
            "author": document.metadata.author,
            "created": str(document.metadata.created_date),
            "page_count": document.metadata.page_count,
            "custom": document.metadata.custom_metadata
        },
        "segments_count": len(document.segments),
        "visual_elements_count": len(document.visual_elements),
        "segments": [],
        "visual_elements": []
    }
    
    # Add all segments with full content
    for i, seg in enumerate(document.segments):
        segment_data = {
            "index": i,
            "type": seg.segment_type,
            "page_number": seg.page_number,
            "content": seg.content,  # Full content
            "metadata": seg.metadata,
            "visual_references": seg.visual_references
        }
        output["segments"].append(segment_data)
    
    # Add visual elements
    for i, ve in enumerate(document.visual_elements):
        visual_data = {
            "index": i,
            "type": ve.element_type.value,
            "page_or_slide": ve.page_or_slide,
            "segment_reference": ve.segment_reference,
            "analysis_metadata": ve.analysis_metadata,
            "file_extension": ve.file_extension
        }
        output["visual_elements"].append(visual_data)
    
    # Save to file
    output_file = f"data/output/{parser_type}_detailed_{timestamp}.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(output, f, indent=2, ensure_ascii=False)
    
    logger.info(f"💾 Detailed output saved to: {output_file}")
    return output_file


async def test_docx_detailed():
    """Test DOCX parser with detailed output"""
    logger.info("\n" + "="*80)
    logger.info("📄 DOCX Parser - Detailed Test")
    logger.info("="*80)
    
    try:
        from core.parsers.implementations.office import DOCXParser
        
        docx_path = Path("data/input/Erkenntnisse_Herausforderungen.docx")
        
        # Create parser
        parser = DOCXParser(config={
            "extract_images": True,
            "extract_tables": True,
            "preserve_formatting": True,
            "table_as_text": True
        }, enable_vlm=False)
        
        # Parse document
        document = await parser.parse(docx_path)
        
        # Log basic info
        logger.info(f"✅ Parsed {len(document.segments)} segments")
        
        # Show all segments
        logger.info("\n--- ALL SEGMENTS ---")
        for i, seg in enumerate(document.segments):
            logger.info(f"\nSegment {i+1}:")
            logger.info(f"  Type: {seg.segment_type}")
            logger.info(f"  Style: {seg.metadata.get('style', 'N/A')}")
            logger.info(f"  Length: {len(seg.content)} chars")
            logger.info(f"  Content: {seg.content}")
            logger.info("-" * 40)
        
        # Save detailed output
        output_file = save_detailed_output(str(docx_path), document, "DOCX")
        
        # Try to extract raw text using python-docx for comparison
        logger.info("\n--- COMPARISON WITH RAW DOCX ---")
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(str(docx_path))
            logger.info(f"Total paragraphs in raw DOCX: {len(doc.paragraphs)}")
            logger.info("\nFirst 5 paragraphs from raw DOCX:")
            for i, para in enumerate(doc.paragraphs[:5]):
                if para.text.strip():
                    logger.info(f"Para {i+1}: {para.text}")
        except Exception as e:
            logger.error(f"Could not compare with raw DOCX: {e}")
        
        return document, output_file
        
    except Exception as e:
        logger.error(f"❌ DOCX test failed: {e}", exc_info=True)
        return None, None


async def test_xlsx_detailed():
    """Test XLSX parser with detailed output"""
    logger.info("\n" + "="*80)
    logger.info("📊 XLSX Parser - Detailed Test")
    logger.info("="*80)
    
    try:
        from core.parsers.implementations.office import XLSXParser
        
        xlsx_path = Path("data/input/PapersWithCode_Abstracts.xlsx")
        
        # Create parser
        parser = XLSXParser(config={
            "extract_charts": True,
            "extract_images": True,
            "extract_formulas": True,
            "sheet_as_segment": True,
            "include_empty_cells": False
        }, enable_vlm=False)
        
        # Parse document
        document = await parser.parse(xlsx_path)
        
        # Log basic info
        logger.info(f"✅ Parsed {len(document.segments)} segments")
        logger.info(f"Sheets: {document.metadata.custom_metadata.get('sheets', [])}")
        
        # Show all segments
        logger.info("\n--- ALL SEGMENTS ---")
        for i, seg in enumerate(document.segments):
            logger.info(f"\nSegment {i+1}:")
            logger.info(f"  Type: {seg.segment_type}")
            logger.info(f"  Sheet: {seg.metadata.get('sheet_name', 'N/A')}")
            logger.info(f"  Length: {len(seg.content)} chars")
            if seg.segment_type == "data_range":
                # Show first few lines for data ranges
                lines = seg.content.split('\n')
                logger.info(f"  First 10 lines:")
                for line in lines[:10]:
                    logger.info(f"    {line}")
                if len(lines) > 10:
                    logger.info(f"    ... ({len(lines) - 10} more lines)")
            else:
                logger.info(f"  Content: {seg.content}")
            logger.info("-" * 40)
        
        # Save detailed output
        output_file = save_detailed_output(str(xlsx_path), document, "XLSX")
        
        # Try to extract raw data using openpyxl for comparison
        logger.info("\n--- COMPARISON WITH RAW XLSX ---")
        try:
            import openpyxl
            wb = openpyxl.load_workbook(xlsx_path, data_only=True)
            logger.info(f"Total sheets in raw XLSX: {len(wb.sheetnames)}")
            logger.info(f"Sheet names: {wb.sheetnames}")
            
            for sheet_name in wb.sheetnames[:1]:  # First sheet only
                sheet = wb[sheet_name]
                logger.info(f"\nSheet '{sheet_name}':")
                logger.info(f"  Dimensions: {sheet.max_row} rows x {sheet.max_column} columns")
                logger.info("  First 5 rows:")
                for row_idx in range(1, min(6, sheet.max_row + 1)):
                    row_data = []
                    for col_idx in range(1, min(10, sheet.max_column + 1)):
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        row_data.append(str(cell.value) if cell.value is not None else "")
                    logger.info(f"    Row {row_idx}: {' | '.join(row_data)}")
        except Exception as e:
            logger.error(f"Could not compare with raw XLSX: {e}")
        
        return document, output_file
        
    except Exception as e:
        logger.error(f"❌ XLSX test failed: {e}", exc_info=True)
        return None, None


async def test_pptx_detailed():
    """Test PPTX parser with detailed output"""
    logger.info("\n" + "="*80)
    logger.info("🎯 PPTX Parser - Detailed Test")
    logger.info("="*80)
    
    try:
        from core.parsers.implementations.office import PPTXParser
        
        pptx_path = Path("data/input/CRISP-DM_TätigkeitenTools.pptx")
        
        # Create parser
        parser = PPTXParser(config={
            "extract_images": True,
            "extract_charts": True,
            "extract_diagrams": True,
            "extract_notes": True,
            "extract_titles": True
        }, enable_vlm=False)
        
        # Parse document
        document = await parser.parse(pptx_path)
        
        # Log basic info
        logger.info(f"✅ Parsed {len(document.segments)} segments")
        logger.info(f"Total slides: {document.metadata.custom_metadata.get('slides', 0)}")
        logger.info(f"Visual elements: {len(document.visual_elements)}")
        
        # Show all segments (slides)
        logger.info("\n--- ALL SLIDES ---")
        for i, seg in enumerate(document.segments):
            if seg.segment_type == "slide":
                logger.info(f"\nSlide {seg.metadata.get('slide_number', i+1)}:")
                logger.info(f"  Title: {seg.metadata.get('slide_title', 'No title')}")
                logger.info(f"  Has notes: {seg.metadata.get('has_notes', False)}")
                logger.info(f"  Content length: {len(seg.content)} chars")
                logger.info(f"  Content:")
                # Show content with indentation
                for line in seg.content.split('\n'):
                    logger.info(f"    {line}")
                logger.info("-" * 40)
        
        # Show visual elements summary
        if document.visual_elements:
            logger.info("\n--- VISUAL ELEMENTS ---")
            visual_summary = {}
            for ve in document.visual_elements:
                ve_type = ve.element_type.value
                visual_summary[ve_type] = visual_summary.get(ve_type, 0) + 1
            logger.info(f"Visual element types: {visual_summary}")
        
        # Save detailed output
        output_file = save_detailed_output(str(pptx_path), document, "PPTX")
        
        # Try to extract raw data using python-pptx for comparison
        logger.info("\n--- COMPARISON WITH RAW PPTX ---")
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))
            logger.info(f"Total slides in raw PPTX: {len(prs.slides)}")
            
            logger.info("\nFirst 3 slides from raw PPTX:")
            for i, slide in enumerate(prs.slides[:3]):
                logger.info(f"\nSlide {i+1}:")
                # Extract text from all shapes
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                logger.info(f"  Text elements: {len(texts)}")
                for j, text in enumerate(texts):
                    logger.info(f"  Text {j+1}: {text[:100]}..." if len(text) > 100 else f"  Text {j+1}: {text}")
        except Exception as e:
            logger.error(f"Could not compare with raw PPTX: {e}")
        
        return document, output_file
        
    except Exception as e:
        logger.error(f"❌ PPTX test failed: {e}", exc_info=True)
        return None, None


async def main():
    """Run detailed tests for all office parsers"""
    logger.info("="*80)
    logger.info("🔍 Office Parser Detailed Content Test")
    logger.info(f"📝 Log file: {log_filename}")
    logger.info("="*80)
    
    # Test each parser
    docx_result = await test_docx_detailed()
    xlsx_result = await test_xlsx_detailed()
    pptx_result = await test_pptx_detailed()
    
    # Summary
    logger.info("\n" + "="*80)
    logger.info("📊 DETAILED TEST SUMMARY")
    logger.info("="*80)
    
    if docx_result[0]:
        logger.info(f"✅ DOCX: Success - Output: {docx_result[1]}")
    else:
        logger.info("❌ DOCX: Failed")
        
    if xlsx_result[0]:
        logger.info(f"✅ XLSX: Success - Output: {xlsx_result[1]}")
    else:
        logger.info("❌ XLSX: Failed")
        
    if pptx_result[0]:
        logger.info(f"✅ PPTX: Success - Output: {pptx_result[1]}")
    else:
        logger.info("❌ PPTX: Failed")
    
    logger.info(f"\n📄 Full log: {log_filename}")
    logger.info("✅ Detailed tests completed!")


if __name__ == "__main__":
    asyncio.run(main())